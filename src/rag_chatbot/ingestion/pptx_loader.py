from pptx import Presentation
import os
import logging


def load_cvs(file_path):
    pptx_files = []

    try:
        pptx_files = list(filter(lambda x: x.endswith(".pptx") , os.listdir(file_path)))
    except Exception:
        logging.warning("could not open path")

    file_data = []

    for file in pptx_files:
        try:
            logging.info(f"opening file {file}")
            presentation = Presentation(file_path + file)
            text_runs = []

            for slide in presentation.slides:
                 for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        runs = []
                        for run in paragraph.runs:
                            runs.append(run.text)

                        text_runs.append("".join(runs))

            text_runs = list(filter(lambda x: x != "", text_runs))

            file_data.append({
                "text": text_runs,
                "file_name": file,
                "candidate_name": f"{text_runs[0]} {text_runs[1]}"
            })

            logging.info(f"file {file} opened succesfully")
            logging.info(f"characters in the file {len("".join(text_runs))}")
        except Exception as ex:
            logging.warning(ex)

    return file_data
